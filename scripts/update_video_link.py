import sys
import pptx
import pathlib
import time
import os
import shutil
import comtypes.client
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Configure UTF-8 encoding
sys.stdout.reconfigure(encoding='utf-8')

VIDEO_URL = "https://drive.google.com/file/d/1s-f2FqNILJzSWFGeqosAJpr7QCAxGOqs/view?usp=drivesdk"

TEXT_WHITE = RGBColor(240, 246, 252)
ACCENT_GREEN = RGBColor(63, 185, 80)
LINK_CYAN = RGBColor(88, 166, 255)

def update_presentation(pptx_file):
    p_path = pathlib.Path(pptx_file).resolve()
    if not p_path.exists():
        print(f"File not found: {pptx_file}")
        return
    prs = pptx.Presentation(p_path)
    slide8 = prs.slides[7] # 0-indexed slide 8
    
    # Locate right card shape
    card_r = slide8.shapes[4]
    for p in card_r.text_frame.paragraphs:
        if "Demonstration Video Link:" in p.text:
            print(f"[{pptx_file}] Updating video link paragraph...")
            p.text = "★ Demonstration Video Link: "
            p.font.name = "Calibri"
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = ACCENT_GREEN
            p.space_before = Pt(7)
            
            run = p.add_run()
            run.text = VIDEO_URL
            run.font.name = "Calibri"
            run.font.size = Pt(10.5)
            run.font.bold = False
            run.font.color.rgb = LINK_CYAN
            run.hyperlink.address = VIDEO_URL
            
    prs.save(str(p_path))
    print(f"[{pptx_file}] PPTX updated successfully.")

def export_pdf(pptx_file, pdf_file):
    src = str(pathlib.Path(pptx_file).resolve())
    pdf = str(pathlib.Path(pdf_file).resolve())
    
    print(f"Exporting PDF from {pptx_file} -> {pdf_file} via PowerPoint COM...")
    try:
        ppt = comtypes.client.CreateObject("PowerPoint.Application")
        ppt.Visible = 1
        pres = ppt.Presentations.Open(src)
        time.sleep(1)
        pres.SaveAs(pdf, 32) # 32 = ppSaveAsPDF
        pres.Close()
        ppt.Quit()
        print(f"SUCCESS: Exported {pdf}")
    except Exception as e:
        print(f"Error during COM export: {e}")

if __name__ == "__main__":
    for f in ["SEMICON_India_Hackathon_SPARTANS_Final.pptx", "SEMICON_India_Hackathon_SemiRestoreAI.pptx"]:
        update_presentation(f)
        
    export_pdf("SEMICON_India_Hackathon_SPARTANS_Final.pptx", "SEMICON_India_Hackathon_SPARTANS_Final.pdf")
    shutil.copy2("SEMICON_India_Hackathon_SPARTANS_Final.pdf", "SEMICON_India_Hackathon_SemiRestoreAI.pdf")
    
    # Also sync to Downloads folder if needed
    dl = pathlib.Path(os.environ["USERPROFILE"]) / "Downloads"
    for name in [
        "SEMICON_India_Hackathon_SPARTANS_Final.pdf",
        "SEMICON_India_Hackathon_SPARTANS_Final.pptx",
        "SEMICON_India_Hackathon_SemiRestoreAI.pdf",
        "SEMICON_India_Hackathon_SemiRestoreAI.pptx"
    ]:
        src = pathlib.Path(name).resolve()
        if src.exists():
            shutil.copy2(str(src), dl / name)
            print(f"Synced to Downloads: {name}")
